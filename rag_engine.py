from __future__ import annotations
from sentence_transformers import SentenceTransformer
from sentence_transformers import CrossEncoder

import math
import os
import re
import string
import unicodedata
from collections import Counter, defaultdict
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

import torch
torch.set_default_device("mps")

from transformers import (
    AutoTokenizer,
    AutoModelForSeq2SeqLM,
)

def extract_text(filepath: str) -> str:
    path = Path(filepath)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(filepath)
    elif suffix == ".docx":
        return _parse_docx(filepath)
    elif suffix == ".pptx":
        return _parse_pptx(filepath)
    elif suffix in (".txt", ".md", ".rst", ".csv", ".json", ".xml", ".html", ".htm"):
        return path.read_text(encoding="utf-8", errors="ignore")
    else:
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            raise ValueError(f"Unsupported file type: {suffix}")


def _parse_pdf(filepath: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(filepath)
    parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            parts.append(f"[Page {i}]\n{text}")
    return "\n\n".join(parts)


def _parse_docx(filepath: str) -> str:
    from docx import Document
    doc = Document(filepath)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def _parse_pptx(filepath: str) -> str:
    from pptx import Presentation
    prs = Presentation(filepath)
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        parts.append("\n".join(slide_parts))
    return "\n\n".join(parts)

@dataclass
class Chunk:
    chunk_id: int
    doc_name: str
    text: str
    start_char: int
    end_char: int
    page_hint: int = 0


def chunk_document(text: str, doc_name: str, chunk_size: int = 350, overlap: int = 80) -> list[Chunk]:
    text = _clean_text(text)
    if not text.strip():
        return []

    sentences = _split_sentences(text)
    if not sentences:
        return []

    chunks: list[Chunk] = []
    buf_sents: list[str] = []
    buf_words = 0
    char_offset = 0
    chunk_start = 0

    for sent in sentences:
        words = sent.split()
        wcount = len(words)

        if buf_words > 0 and buf_words + wcount > chunk_size:
            _flush(chunks, buf_sents, doc_name, chunk_start, char_offset)
            buf_sents, buf_words = _trim_overlap(buf_sents, overlap)
            chunk_start = char_offset - sum(len(s) + 1 for s in buf_sents)

        buf_sents.append(sent)
        buf_words += wcount
        char_offset += len(sent) + 1

    if buf_sents:
        _flush(chunks, buf_sents, doc_name, chunk_start, char_offset)

    for i, c in enumerate(chunks):
        c.chunk_id = i

    return chunks


def _flush(
    chunks: list,
    sents: list,
    doc_name: str,
    start: int,
    end: int
):
    text = " ".join(sents).strip()

    if not text:
        return

    # Remove noisy chunks
    noise_ratio = (
        len(re.findall(r"[^\w\s]", text))
        / max(1, len(text))
    )

    if noise_ratio > 0.35:
        return

    if len(text.split()) < 10:
        return

    chunks.append(
        Chunk(
            chunk_id=len(chunks),
            doc_name=doc_name,
            text=text,
            start_char=start,
            end_char=end,
            page_hint=_estimate_page(start),
        )
    )


def _trim_overlap(sents: list[str], target_words: int) -> tuple[list[str], int]:
    result, wc = [], 0
    for s in reversed(sents):
        w = len(s.split())
        if wc + w > target_words and result:
            break
        result.insert(0, s)
        wc += w
    return result, wc


def _clean_text(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def _split_sentences(text: str) -> list[str]:
    text = text.replace("\n", " ")
    sentences = re.split(r'(?<=[.!?])\s+', text)
    cleaned = []
    for s in sentences:
        s = s.strip()
        if len(s.split()) >= 5:
            cleaned.append(s)

    return cleaned


def _estimate_page(char_pos: int) -> int:
    return max(1, char_pos // 3000 + 1)


class BM25:
    def __init__(self, k1: float = 1.5, b: float = 0.75):
        self.k1, self.b = k1, b
        self.corpus: list[list[str]] = []
        self.doc_freqs: list[Counter] = []
        self.idf: dict[str, float] = {}
        self.avgdl: float = 0.0
        self.N: int = 0

    def fit(self, texts: list[str]):
        self.corpus = [self._tok(t) for t in texts]
        self.N = len(self.corpus)
        self.avgdl = sum(len(d) for d in self.corpus) / max(1, self.N)
        self.doc_freqs = [Counter(d) for d in self.corpus]
        df: Counter = Counter()
        for doc in self.corpus:
            for term in set(doc):
                df[term] += 1
        self.idf = {
            term: math.log((self.N - f + 0.5) / (f + 0.5) + 1)
            for term, f in df.items()
        }

    def score(self, query: str) -> np.ndarray:
        tokens = self._tok(query)
        scores = np.zeros(self.N)
        for tok in tokens:
            if tok not in self.idf:
                continue
            idf = self.idf[tok]
            for i, doc in enumerate(self.corpus):
                tf = self.doc_freqs[i].get(tok, 0)
                dl = len(doc)
                denom = tf + self.k1 * (1 - self.b + self.b * dl / max(1, self.avgdl))
                scores[i] += idf * tf * (self.k1 + 1) / denom
        return scores

    @staticmethod
    def _tok(text: str) -> list[str]:
        return text.lower().translate(str.maketrans("", "", string.punctuation)).split()


@dataclass
class RetrievalResult:
    chunk: Chunk
    score: float
    tfidf_score: float
    bm25_score: float


class HybridRetriever:
    def __init__(self, alpha: float = 0.55):
        self.reranker = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")
        self.alpha = alpha
        self.chunks: list[Chunk] = []
        self.tfidf = TfidfVectorizer(
            ngram_range=(1, 2),
            sublinear_tf=True,
            strip_accents="unicode",
            min_df=1,
        )
        self.bm25 = BM25()
        self.tfidf_matrix = None
        self._fitted = False

        self.embedder = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")
        self.chunk_embeddings = None

    def index(self, chunks: list[Chunk]):
        self.chunks = chunks

        if not chunks:
            self._fitted = False
            return

        texts = [c.text for c in chunks]
        self.tfidf_matrix = self.tfidf.fit_transform(texts)
        self.bm25.fit(texts)
        self.chunk_embeddings = self.embedder.encode(
            texts,
            convert_to_numpy=True,
            normalize_embeddings=True,
            show_progress_bar=False,
        )
        self._fitted = True

    def retrieve(
        self,
        query: str,
        top_k: int = 6,
        doc_filter: list[str] | None = None,
    ) -> list[RetrievalResult]:

        if not self._fitted or not self.chunks:
            return []

        q_vec = self.tfidf.transform([query])
        tfidf_scores = cosine_similarity(
            q_vec,
            self.tfidf_matrix
        ).flatten()

        bm25_scores = self.bm25.score(query)
        query_embedding = self.embedder.encode(
            [query],
            convert_to_numpy=True,
            normalize_embeddings=True,
        )

        semantic_scores = cosine_similarity(
            query_embedding,
            self.chunk_embeddings
        ).flatten()

        tfidf_norm = _normalize(tfidf_scores)
        bm25_norm = _normalize(bm25_scores)
        semantic_norm = _normalize(semantic_scores)

        combined = (
            semantic_norm * 0.60 +
            bm25_norm * 0.25 +
            tfidf_norm * 0.15
        )

        if doc_filter:
            filter_set = set(doc_filter)

            for i, chunk in enumerate(self.chunks):
                if chunk.doc_name not in filter_set:
                    combined[i] = -1

        candidate_pool = np.argsort(combined)[::-1][:50]

        pool_embeddings = self.chunk_embeddings[candidate_pool]
        pool_scores = combined[candidate_pool]

        selected_local = mmr_select(
            pool_embeddings,
            pool_scores,
            top_k=top_k * 5,
        )

        top_indices = [
            candidate_pool[i]
            for i in selected_local
        ]
        pairs = [
            (query, self.chunks[idx].text)
            for idx in top_indices
        ]

        rerank_scores = self.reranker.predict(pairs)
        reranked = sorted(
            zip(top_indices, rerank_scores),
            key=lambda x: x[1],
            reverse=True
        )
        results = []

        for idx, rerank_score in reranked[:top_k]:
            results.append(
                RetrievalResult(
                    chunk=self.chunks[idx],
                    score=float(rerank_score),
                    tfidf_score=float(tfidf_norm[idx]),
                    bm25_score=float(bm25_norm[idx]),
                )
            )

        return results[:top_k]


def _normalize(arr: np.ndarray) -> np.ndarray:
    mn, mx = arr.min(), arr.max()
    if mx - mn < 1e-10:
        return np.zeros_like(arr, dtype=float)
    return (arr - mn) / (mx - mn)

STOPWORDS = {
    "a","an","the","is","are","was","were","be","been","being","have","has","had",
    "do","does","did","will","would","could","should","may","might","shall","can",
    "to","of","in","for","on","with","at","by","from","as","into","through","during",
    "before","after","above","below","and","but","or","nor","so","yet","both","either",
    "not","no","what","which","who","whom","how","when","where","why","that","this",
    "these","those","i","me","my","we","our","you","your","he","his","she","her","it",
    "its","they","their","them","all","any","each","few","more","most","other","some",
    "such","own","same","than","too","very","just","about","up","out","if","then","s",
}


def generate_answer(
    query: str,
    results: list[RetrievalResult]
) -> dict[str, Any]:

    if not results:
        return {
            "answer": "No relevant information found.",
            "confidence": 0,
            "sources": [],
            "chunks_used": 0,
        }

    query_terms = _extract_terms(query)

    selected_sentences = []
    seen = set()

    for result in results:
        sentences = _split_sentences(result.chunk.text)

        for sent in sentences:
            sent_clean = sent.strip()

            if not sent_clean:
                continue

            sent_words = set(
                re.findall(
                    r"\b[a-zA-Z0-9]+\b",
                    sent_clean.lower()
                )
            )

            # Query coverage
            overlap = len(
                sent_words.intersection(query_terms)
            )

            coverage = overlap / max(1, len(query_terms))

            # Semantic weighting
            score = (
                coverage * 0.7 +
                result.score * 0.3
            )

            # Ignore weak matches
            if coverage < 0.15:
                continue

            # Remove duplicate-ish sentences
            fingerprint = sent_clean.lower()[:120]

            if fingerprint in seen:
                continue

            seen.add(fingerprint)

            selected_sentences.append({
                "text": sent_clean,
                "score": score,
                "doc": result.chunk.doc_name,
            })

    # Sort best first
    selected_sentences.sort(
        key=lambda x: x["score"],
        reverse=True
    )

    # Dynamic answer assembly
    answer_parts = []

    used_words = 0
    max_words = 250

    for item in selected_sentences:
        wc = len(item["text"].split())

        if used_words + wc > max_words:
            break

        answer_parts.append(item["text"])
        used_words += wc

    if not answer_parts:
        # fallback
        answer_parts = [
            results[0].chunk.text[:500]
        ]

    answer = "\n\n".join(answer_parts)

    # Sources
    sources = []

    used_docs = set()

    for r in results:
        if r.chunk.doc_name in used_docs:
            continue

        used_docs.add(r.chunk.doc_name)

        sources.append({
            "document": r.chunk.doc_name,
            "relevance": round(r.score * 100, 1),
            "page_hint": r.chunk.page_hint,
            "excerpt": r.chunk.text[:300],
        })

    confidence = np.mean([
        r.score for r in results[:3]
    ])
    
    # Simple sigmoid-like normalization for CrossEncoder scores (which often range from -10 to +10)
    normalized_confidence = 1 / (1 + np.exp(-confidence))

    return {
        "answer": answer,
        "confidence": round(float(normalized_confidence) * 100, 1),
        "sources": sources,
        "chunks_used": len(results),
    }


def mmr_select(
    embeddings,
    scores,
    top_k=6,
    lambda_param=0.7
):
    selected = []

    candidates = list(range(len(scores)))

    while len(selected) < min(top_k, len(scores)):

        if not selected:
            idx = int(np.argmax(scores))
            selected.append(idx)
            candidates.remove(idx)
            continue

        mmr_scores = []

        for c in candidates:
            relevance = scores[c]

            diversity = max([
                cosine_similarity(
                    [embeddings[c]],
                    [embeddings[s]]
                )[0][0]
                for s in selected
            ])

            mmr = (
                lambda_param * relevance -
                (1 - lambda_param) * diversity
            )

            mmr_scores.append((c, mmr))

        best = max(mmr_scores, key=lambda x: x[1])[0]

        selected.append(best)
        candidates.remove(best)

    return selected

def _extract_terms(query: str) -> set[str]:
    words = re.findall(r"\b[a-zA-Z0-9]+\b", query.lower())
    return {w for w in words if w not in STOPWORDS and len(w) > 2}


def _score_sentence(sentence: str, query_terms: set[str]) -> float:
    words = re.findall(r"\b[a-zA-Z0-9]+\b", sentence.lower())
    if not words:
        return 0.0
    hits = sum(1 for w in words if w in query_terms)
    tf = hits / len(words)

    length_factor = min(1.0, len(words) / 20.0) * (1.0 if len(words) <= 60 else 0.6)
    return tf * 0.75 + length_factor * 0.25


def _build_answer_text(query: str, sentences: list[str], query_terms: set[str]) -> str:
    if not sentences:
        return "The documents don't contain a clear answer to this question."

    highlighted = [_highlight(s, query_terms) for s in sentences]
    
    return "\n\n".join(highlighted)


def _highlight(text: str, terms: set[str]) -> str:
    def repl(m: re.Match) -> str:
        w = m.group(0)
        if w.lower() in terms:
            return f"<mark>{w}</mark>"
        return w
    return re.sub(r"\b[a-zA-Z0-9]+\b", repl, text)

class DocumentStore:
    def __init__(self):
        MODEL_NAME = "google/flan-t5-base"
        self.llm_tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME)
        self.llm_model = AutoModelForSeq2SeqLM.from_pretrained (
            MODEL_NAME,
            dtype=torch.float16,
            device_map="auto",
        )
        self.documents: dict[str, dict] = {}
        self.all_chunks: list[Chunk] = []
        self.retriever = HybridRetriever(alpha=0.55)

    def generate_llm_answer(
        self,
        question: str,
        context: str,
    ) -> str:
        context = context[:5000]
        prompt = f"""
TASK: Extract all relevant information from the context to answer the question.
FORMAT: Use a bulleted list for multiple items.
IMPORTANT: Include ALL names, companies, or steps mentioned. Do not summarize or skip anything.

CONTEXT:
{context}

QUESTION: {question}

ANSWER:
"""
        inputs = self.llm_tokenizer(
            prompt,
            return_tensors="pt",
            truncation=True,
            max_length=2048,
        ).to(self.llm_model.device)

        outputs = self.llm_model.generate(
            **inputs,
            max_new_tokens=450,
            temperature=0.1,
            do_sample=True,
            top_p=0.9,
            repetition_penalty=1.5,
            no_repeat_ngram_size=3,
        )

        answer = self.llm_tokenizer.decode(
            outputs[0],
            skip_special_tokens=True,
        )

        return answer.strip()

    def add_document(self, filepath: str, original_name: str) -> dict:
        try:
            text = extract_text(filepath)
        except Exception as e:
            return {"success": False, "error": str(e)}

        if not text.strip():
            return {"success": False, "error": "File appears to be empty or could not be parsed."}

        chunks = chunk_document(text, original_name)
        if not chunks:
            return {"success": False, "error": "No content could be extracted from this file."}

        # Remove old chunks for this doc
        self.all_chunks = [c for c in self.all_chunks if c.doc_name != original_name]
        # Renumber and append
        offset = len(self.all_chunks)
        for i, c in enumerate(chunks):
            c.chunk_id = offset + i
        self.all_chunks.extend(chunks)

        self.documents[original_name] = {
            "name": original_name,
            "path": filepath,
            "char_count": len(text),
            "chunk_count": len(chunks),
            "word_count": len(text.split()),
        }

        self._reindex()
        return {
            "success": True,
            "name": original_name,
            "chunks": len(chunks),
            "words": self.documents[original_name]["word_count"],
        }

    def remove_document(self, doc_name: str):
        self.documents.pop(doc_name, None)
        self.all_chunks = [c for c in self.all_chunks if c.doc_name != doc_name]
        for i, c in enumerate(self.all_chunks):
            c.chunk_id = i
        self._reindex()

    def _reindex(self):
        self.retriever.index(self.all_chunks)

    def query(self, question: str, top_k: int = 10, doc_filter: list[str] | None = None) -> dict:
        if not self.all_chunks:
            return {
                "answer": "No documents uploaded yet. Please upload files first.",
                "confidence": 0.0, "sources": [], "chunks_used": 0,
            }
        results = self.retriever.retrieve(question, top_k=top_k, doc_filter=doc_filter)
        retrieved_text = "\n\n".join([
            r.chunk.text
            for r in results
        ])

        llm_answer = self.generate_llm_answer(
            question,
            retrieved_text,
        )
        base = generate_answer(question, results)
        base["answer"] = llm_answer

        return base

    def list_documents(self) -> list[dict]:
        return list(self.documents.values())

    def clear(self):
        self.documents.clear()
        self.all_chunks.clear()
        self.retriever = HybridRetriever(alpha=0.55)
